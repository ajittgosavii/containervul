"""Generate the ContainerVul product presentation PPTX."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
DARK_BG = RGBColor(0x0F, 0x34, 0x60)
PURPLE = RGBColor(0x53, 0x34, 0x83)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF7, 0xFF)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_RED = RGBColor(0xF5, 0x57, 0x6C)
ACCENT_BLUE = RGBColor(0x4F, 0xAC, 0xFE)
ACCENT_GREEN = RGBColor(0x43, 0xE9, 0x7B)
ACCENT_ORANGE = RGBColor(0xFF, 0x99, 0x00)
ACCENT_PURPLE = RGBColor(0xA1, 0x8C, 0xD1)
GRAY_TEXT = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, corner_radius=Inches(0.15)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_bullet_text(slide, left, top, width, height, bullets, font_size=14, color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.level = 0
    return txBox

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, DARK_BG)

# Accent bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
    "Enterprise Container Vulnerability\nManagement Platform", 44, WHITE, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
    "Multi-Cloud Security  •  Agentic AI  •  MCP Integration  •  ServiceNow  •  Compliance", 20, RGBColor(0xBB, 0xBB, 0xDD), False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
    "AWS EKS/ECS  •  Azure AKS/ACI  •  GCP GKE/Cloud Run  •  Multi-Account  •  CIS & NIST Compliance", 16, RGBColor(0x99, 0x99, 0xBB), False, PP_ALIGN.CENTER)

# Bottom bar with version
add_shape(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.7), PURPLE)
add_text_box(slide, Inches(1), Inches(6.9), Inches(11), Inches(0.4),
    "ContainerVul v2.0  |  Powered by Claude AI  |  Open Source", 14, WHITE, False, PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 2: Problem Statement
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6), "The Container Security Challenge", 32, WHITE, True)

problems = [
    ("85%", "of containers have at least one\npatchable vulnerability", ACCENT_RED),
    ("60%", "of organizations lack visibility\nacross multi-cloud containers", ACCENT_ORANGE),
    ("45 days", "average time to remediate\na critical container CVE", ACCENT_PURPLE),
    ("3x", "cloud providers with different\nsecurity tools and APIs", ACCENT_BLUE),
]

for i, (stat, desc, color) in enumerate(problems):
    x = Inches(0.6 + i * 3.1)
    card = add_shape(slide, x, Inches(1.6), Inches(2.8), Inches(2.5), LIGHT_GRAY)
    add_shape(slide, x, Inches(1.6), Inches(2.8), Inches(0.08), color)
    add_text_box(slide, x + Inches(0.3), Inches(1.9), Inches(2.2), Inches(0.7), stat, 36, color, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.7), Inches(2.4), Inches(0.8), desc, 14, GRAY_TEXT, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(1.5),
    "Organizations need a unified platform that scans containers across all cloud providers,\n"
    "leverages AI for intelligent remediation, integrates with ITSM workflows,\n"
    "and ensures compliance — all from a single pane of glass.",
    16, DARK_TEXT, False, PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 3: Solution Overview
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6), "ContainerVul — Solution Overview", 32, WHITE, True)

features = [
    ("Multi-Cloud Scanning", "EKS, ECS, AKS, ACI, GKE,\nCloud Run, ECR, ACR, GAR\n9 container services", ACCENT_BLUE),
    ("Agentic AI", "Claude tool-use agent with\n13 tools for autonomous\nscan → analyze → remediate", ACCENT_PURPLE),
    ("MCP Server", "15 MCP tools exposing all\ncapabilities to Claude Desktop\nand other MCP clients", RGBColor(0x93, 0x33, 0xEA)),
    ("ServiceNow", "Auto-create P1/P2 incidents\nCMDB sync, change requests\nBidirectional status sync", ACCENT_GREEN),
    ("Compliance", "CIS Docker Benchmark\nCIS Kubernetes Benchmark\nNIST SP 800-190", ACCENT_RED),
    ("Enterprise", "RBAC (4 roles), audit logging\nMulti-account support\nMulti-format reports", ACCENT_ORANGE),
]

for i, (title, desc, color) in enumerate(features):
    row = i // 3
    col = i % 3
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.5 + row * 2.6)
    card = add_shape(slide, x, y, Inches(3.8), Inches(2.2), LIGHT_GRAY)
    header = add_shape(slide, x, y, Inches(3.8), Inches(0.5), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.05), Inches(3.4), Inches(0.4), title, 18, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(1.3), desc, 14, DARK_TEXT, False, PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 4: Architecture
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6), "Platform Architecture", 32, WHITE, True)

# Layer boxes
layers = [
    ("Users & Interfaces", "Browser  •  Streamlit UI  •  AI Agent Chat  •  MCP Clients  •  Claude Desktop  •  REST API", RGBColor(0x66, 0x7E, 0xEA), Inches(1.3)),
    ("Application Layer", "Core Scanning  •  Agentic AI (13 tools)  •  Multi-Cloud Module  •  Enterprise & Integrations", ACCENT_RED, Inches(2.4)),
    ("External Services", "AWS (EKS/ECS/ECR)  •  Azure (AKS/ACI/ACR)  •  GCP (GKE/Run/GAR)  •  NIST NVD  •  Claude API  •  ServiceNow  •  Trivy", ACCENT_BLUE, Inches(3.5)),
    ("Data Layer", "Session State  •  SQLite  •  TTL Cache  •  Audit Log  •  ServiceNow Ticket Cache", ACCENT_ORANGE, Inches(4.6)),
]

for title, desc, color, y in layers:
    add_shape(slide, Inches(0.6), y, Inches(12.1), Inches(0.9), color)
    add_text_box(slide, Inches(0.9), y + Inches(0.05), Inches(3), Inches(0.35), title, 16, WHITE, True)
    add_text_box(slide, Inches(0.9), y + Inches(0.4), Inches(11), Inches(0.4), desc, 12, WHITE)

# Key metrics
metrics = [("85+", "Python Files"), ("13", "AI Agent Tools"), ("15", "MCP Tools"), ("9", "Container\nServices"), ("3", "Compliance\nFrameworks"), ("3", "Cloud\nProviders")]
for i, (num, label) in enumerate(metrics):
    x = Inches(0.6 + i * 2.1)
    add_shape(slide, x, Inches(5.8), Inches(1.8), Inches(1.2), DARK_BG)
    add_text_box(slide, x, Inches(5.9), Inches(1.8), Inches(0.6), num, 30, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(6.4), Inches(1.8), Inches(0.5), label, 11, RGBColor(0xBB, 0xBB, 0xDD), False, PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 5: Agentic AI Deep Dive
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), PURPLE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6), "Agentic AI — Claude Tool-Use Architecture", 32, WHITE, True)

# Agent loop
add_shape(slide, Inches(0.5), Inches(1.4), Inches(5.8), Inches(5.5), LIGHT_GRAY)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.4), "Agent Loop (max 15 turns)", 18, DARK_TEXT, True)

steps = [
    "1. User sends natural language query",
    "2. Claude analyzes intent, selects tools",
    "3. ToolExecutor dispatches to handler",
    "4. Handler calls core/cloud/snow module",
    "5. Result returned to Claude as tool_result",
    "6. Claude reasons over results",
    "7. May call more tools or respond to user",
    "8. Discovered vulns synced to session state",
]
add_bullet_text(slide, Inches(0.8), Inches(2.1), Inches(5), Inches(4), steps, 14, DARK_TEXT)

# Tools list
add_shape(slide, Inches(6.8), Inches(1.4), Inches(6), Inches(5.5), LIGHT_GRAY)
add_text_box(slide, Inches(7.1), Inches(1.5), Inches(5), Inches(0.4), "13 Agent Tools", 18, DARK_TEXT, True)

tools_left = [
    "scan_dockerfile — Regex pattern analysis",
    "lookup_cve — NIST NVD API query",
    "search_product_cves — Product search",
    "calculate_risk_score — Risk assessment",
    "generate_remediation_plan — AI plan",
    "scan_cloud_service — EKS/AKS/GKE...",
    "check_compliance — CIS / NIST",
]
tools_right = [
    "list_cloud_accounts — Account mgmt",
    "servicenow_create_incident — P1-P4",
    "servicenow_bulk_create — Mass tickets",
    "servicenow_search_tickets — Query",
    "servicenow_sync_cmdb — CI registration",
    "servicenow_create_change — CHG request",
]
add_bullet_text(slide, Inches(7.1), Inches(2.1), Inches(5.5), Inches(2.5), tools_left, 12, DARK_TEXT)
add_bullet_text(slide, Inches(7.1), Inches(4.5), Inches(5.5), Inches(2), tools_right, 12, ACCENT_GREEN)

# ============================================================================
# SLIDE 6: ServiceNow Integration
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), ACCENT_GREEN)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6), "ServiceNow Integration", 32, WHITE, True)

snow_features = [
    ("Incident Management", [
        "Auto-create P1/P2 incidents for CRITICAL/HIGH vulns",
        "Severity → Priority mapping (configurable)",
        "Structured descriptions with CVE, CVSS, remediation",
        "Bulk creation with deduplication",
        "Bidirectional status sync",
    ], ACCENT_RED),
    ("CMDB Sync", [
        "Container images → CI records",
        "Kubernetes clusters → CI records",
        "Services with parent-child relationships",
        "Operational status based on vuln severity",
        "Auto-discovery from cloud scans",
    ], ACCENT_BLUE),
    ("Change Management", [
        "Change requests for remediation actions",
        "Auto-classify: standard vs normal",
        "Implementation & backout plans populated",
        "Risk assessment included",
        "Links to vulnerability tickets",
    ], ACCENT_PURPLE),
]

for i, (title, bullets, color) in enumerate(snow_features):
    x = Inches(0.5 + i * 4.2)
    add_shape(slide, x, Inches(1.4), Inches(3.9), Inches(4.2), LIGHT_GRAY)
    add_shape(slide, x, Inches(1.4), Inches(3.9), Inches(0.5), color)
    add_text_box(slide, x + Inches(0.2), Inches(1.45), Inches(3.5), Inches(0.4), title, 16, WHITE, True, PP_ALIGN.CENTER)
    add_bullet_text(slide, x + Inches(0.3), Inches(2.1), Inches(3.3), Inches(3.2), bullets, 12, DARK_TEXT)

add_text_box(slide, Inches(0.5), Inches(5.9), Inches(12), Inches(0.8),
    "Supports Basic Auth & OAuth2  •  Retry with exponential backoff  •  Rate limit handling  •  Connection pooling",
    14, GRAY_TEXT, False, PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 7: Multi-Cloud Support
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), ACCENT_BLUE)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6), "Multi-Cloud Container Support", 32, WHITE, True)

clouds = [
    ("Amazon Web Services", [
        "EKS — Kubernetes cluster scanning",
        "ECS — Task definitions & Fargate",
        "ECR — Native image scan findings",
        "STS — Cross-account assume-role",
        "Multi-region support",
    ], ACCENT_ORANGE),
    ("Microsoft Azure", [
        "AKS — Managed Kubernetes",
        "ACI — Container Instances",
        "ACR — Container Registry",
        "Multi-subscription support",
        "Service Principal & DefaultCredential",
    ], RGBColor(0x00, 0x78, 0xD4)),
    ("Google Cloud Platform", [
        "GKE — Google Kubernetes Engine",
        "Cloud Run — Serverless containers",
        "Artifact Registry — Image storage",
        "Container Analysis API — Vuln scan",
        "Multi-project support",
    ], RGBColor(0x42, 0x85, 0xF4)),
]

for i, (title, bullets, color) in enumerate(clouds):
    x = Inches(0.5 + i * 4.2)
    add_shape(slide, x, Inches(1.4), Inches(3.9), Inches(4.5), LIGHT_GRAY)
    add_shape(slide, x, Inches(1.4), Inches(3.9), Inches(0.55), color)
    add_text_box(slide, x + Inches(0.2), Inches(1.47), Inches(3.5), Inches(0.4), title, 16, WHITE, True, PP_ALIGN.CENTER)
    add_bullet_text(slide, x + Inches(0.3), Inches(2.2), Inches(3.3), Inches(3.5), bullets, 13, DARK_TEXT)

# ============================================================================
# SLIDE 8: Compliance Frameworks
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), ACCENT_RED)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6), "Compliance & Governance", 32, WHITE, True)

frameworks = [
    ("CIS Docker Benchmark v1.6", "9 Controls", [
        "4.1 — Non-root user for containers",
        "4.2 — Trusted base images only",
        "4.6 — HEALTHCHECK instructions",
        "4.9 — COPY instead of ADD",
        "4.10 — No secrets in Dockerfiles",
        "5.8 — No privileged mode",
    ]),
    ("CIS Kubernetes Benchmark", "8 Controls", [
        "5.1.1 — Restrict cluster-admin",
        "5.2.1 — Minimize privileged containers",
        "5.2.2 — No root containers",
        "5.2.6 — Minimize capabilities",
        "5.4.1 — Secrets as files, not ENV",
        "5.7.3 — Security context enforcement",
    ]),
    ("NIST SP 800-190", "8 Controls", [
        "3.1.1 — Image vulnerabilities",
        "3.1.2 — Image configuration defects",
        "3.1.4 — Embedded secrets",
        "3.3.1 — Network access controls",
        "3.4.1 — Runtime vulnerabilities",
        "3.5.1 — Host OS hardening",
    ]),
]

for i, (title, count, controls) in enumerate(frameworks):
    x = Inches(0.5 + i * 4.2)
    add_shape(slide, x, Inches(1.4), Inches(3.9), Inches(5.2), LIGHT_GRAY)
    add_shape(slide, x, Inches(1.4), Inches(3.9), Inches(0.7), DARK_BG)
    add_text_box(slide, x + Inches(0.2), Inches(1.45), Inches(3.5), Inches(0.35), title, 15, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(1.78), Inches(3.5), Inches(0.25), count, 12, RGBColor(0x99, 0x99, 0xBB), False, PP_ALIGN.CENTER)
    add_bullet_text(slide, x + Inches(0.3), Inches(2.4), Inches(3.3), Inches(4), controls, 12, DARK_TEXT)

# ============================================================================
# SLIDE 9: Technology Stack
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), DARK_BG)
add_text_box(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6), "Technology Stack", 32, WHITE, True)

stack = [
    ("Frontend", "Streamlit, Plotly, Custom CSS\nInteractive charts, responsive layout\n11 modular pages", RGBColor(0x66, 0x7E, 0xEA)),
    ("AI / ML", "Anthropic Claude API (tool_use)\nFastMCP Server (stdio/SSE)\n13 agent tools, system prompts", ACCENT_PURPLE),
    ("Cloud SDKs", "boto3 (AWS), azure-identity\ngoogle-cloud-container\nkubernetes client", ACCENT_BLUE),
    ("Data", "Pydantic v2 (20+ models)\npydantic-settings, cachetools\nSQLite (optional), session state", ACCENT_ORANGE),
    ("Security", "NIST NVD API, Trivy CLI\nRegex pattern engine\nCIS/NIST compliance checks", ACCENT_RED),
    ("Integration", "ServiceNow REST API\nBasic auth + OAuth2\nCMDB, Incidents, Changes", ACCENT_GREEN),
]

for i, (title, desc, color) in enumerate(stack):
    row = i // 3
    col = i % 3
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.4 + row * 2.8)
    add_shape(slide, x, y, Inches(3.8), Inches(2.4), LIGHT_GRAY)
    add_shape(slide, x, y, Inches(3.8), Inches(0.5), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.05), Inches(3.4), Inches(0.4), title, 18, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.2), Inches(1.5), desc, 14, DARK_TEXT, False, PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 10: Getting Started / CTA
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), PURPLE)

add_text_box(slide, Inches(1), Inches(1), Inches(11), Inches(0.8),
    "Get Started", 40, WHITE, True, PP_ALIGN.CENTER)

steps = [
    ("1", "Install", "pip install -r requirements.txt", ACCENT_BLUE),
    ("2", "Configure", "Set CLAUDE_API_KEY, cloud\ncredentials, ServiceNow", ACCENT_GREEN),
    ("3", "Launch", "streamlit run streamlit_app.py", ACCENT_PURPLE),
    ("4", "Scan", "Analyze Dockerfiles, scan\ncloud services, chat with AI", ACCENT_RED),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = Inches(0.8 + i * 3.1)
    add_shape(slide, x, Inches(2.2), Inches(2.8), Inches(2.5), color)
    add_text_box(slide, x, Inches(2.3), Inches(2.8), Inches(0.6), num, 36, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.9), Inches(2.8), Inches(0.4), title, 20, WHITE, True, PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.5), Inches(2.4), Inches(1), desc, 13, WHITE, False, PP_ALIGN.CENTER)

# MCP entry point
add_shape(slide, Inches(2), Inches(5.2), Inches(9.3), Inches(0.7), PURPLE)
add_text_box(slide, Inches(2), Inches(5.25), Inches(9.3), Inches(0.6),
    "MCP Server:  python -m containervul.mcp.server   |   containervul-mcp", 16, WHITE, False, PP_ALIGN.CENTER)

# Footer
add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.5),
    "github.com/ajittgosavii/containervul", 18, ACCENT_BLUE, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(6.7), Inches(11), Inches(0.4),
    "ContainerVul v2.0  •  Open Source  •  Powered by Claude AI & MCP", 13, RGBColor(0x99, 0x99, 0xBB), False, PP_ALIGN.CENTER)

# Save
out_path = os.path.join(os.path.dirname(__file__), "ContainerVul_Product_Presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
